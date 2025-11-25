from fastapi import FastAPI, File, UploadFile, HTTPException,Query
from fastapi.responses import FileResponse,StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from langchain_text_splitters  import RecursiveCharacterTextSplitter   
from langchain_google_genai import GoogleGenerativeAIEmbeddings,ChatGoogleGenerativeAI
from pinecone import Pinecone
from pypdf import PdfReader
import uuid
from pdf2docx import Converter
import os
from pydantic import BaseModel
from pathlib import Path
import psycopg2
import uuid
from typing import Optional
from dotenv import load_dotenv
import fitz
import unicodedata
import os
import io
from html2docx import html2docx
from xhtml2pdf import pisa
from pptx import Presentation
load_dotenv()



app = FastAPI()
conn = psycopg2.connect(
    host="localhost",
    database="RAG",
    user="postgres",
    password="postgres"
)

cursor = conn.cursor()
pc=Pinecone(api_key=os.getenv("api_key_rag"))
index = pc.Index("pdf-rag-index")
embeddings = GoogleGenerativeAIEmbeddings(model="text-embedding-004")
llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash",api_key=os.getenv("GOOGLE_API_KEY"),temperature=0.2)



app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


UPLOAD_FOLDER = os.path.join(os.getcwd(), "uploads")
class QueryRequest(BaseModel):
    question: str
    context: str | None = None


os.makedirs(UPLOAD_FOLDER, exist_ok=True) 

ALLOWED_EXTENSIONS = {'pdf'}
MAX_FILE_SIZE = 32 * 1024 * 1024  

@app.post("/export-docx")
async def export_docx(data: dict):
   
    html = data["html"]
  
    
    


    file_stream = html2docx(html, "My Document")
    file_stream.seek(0)
    

    return StreamingResponse(
        file_stream,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": "attachment; filename=document.docx"},
    )
    
@app.post("/export-pdf")
async def export_docx(data: dict):
   
    html = data["html"]
  
    pdf_stream = io.BytesIO()
    pisa.CreatePDF(
        html,
        dest=pdf_stream
    )


    pdf_stream.seek(0)
    

    return StreamingResponse(
        pdf_stream,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=document.pdf"},
    )
@app.post("/resolve")
async def resolve_query(req:QueryRequest):
    try:
        question = req.question.strip()
        context=req.context.strip()
       
        
       
      
        prompt = f"""
        You are an expert AI assistant answering questions based on your own knowledge:

        CONTEXT:
        {context}

        INSTRUCTION:
        {question} 

        ANSWER:
        """
        answer = llm.invoke(prompt)
        return {
            "question": question,
            "answer": answer,
         
        }
        
    except Exception as e:
        raise e
@app.post("/query")
async def query(req:QueryRequest,userid:str=Query(...),conversation_id:str = Query(...)):
    try:
        question = req.question.strip()
        save_message(conversation_id=conversation_id,role="user",content=question)        
        embed_ques=embeddings.embed_query(question)
        results = index.query(
            vector=embed_ques,
            top_k=3,
            filter={"userId": userid},
            include_metadata=True
          
        )

      
        context = ""
        for match in results["matches"]:
            context += match["metadata"]["text"] + "\n"
     
        prompt = f"""
        You are an expert AI assistant answering questions based on your own knowledge and the following PDF context:

        CONTEXT:
        {context}

        QUESTION:
        {question} 

        ANSWER:
        """
        answer = llm.invoke(prompt)
      

        save_message(conversation_id=conversation_id,role="assistant",content=answer.content)
        
        return {
            "question": question,
            "answer": answer,
         
        }

    except Exception as e:
        print(e)
        
def upload_pdf(file: UploadFile = File(...), userid: str = Query(...)):
    try:
        conversation_id = str(uuid.uuid4())
        cursor.execute(
            "INSERT INTO conversations (id, user_id, title) VALUES (%s, %s, %s)",
            (conversation_id, userid, file.filename)
        )
        conn.commit()
        return {"conversation_id": conversation_id}
    except Exception as e:
        raise e

def save_message(conversation_id: str, role: str, content: str):
    try:
        cursor.execute(
            "INSERT INTO messages (conversation_id, role, content) VALUES (%s, %s, %s)",
            (conversation_id, role, content)
        )
        conn.commit()
    except Exception as e:
        raise e     
    
    
@app.get("/messages")
async def get_messages(conversation_id: str = Query(...)):
    cursor.execute("""
        SELECT id, role, content, created_at
        FROM messages
        WHERE conversation_id = %s
        ORDER BY created_at ASC
    """, (conversation_id,))
    
    rows = cursor.fetchall()
    
    
    return {
        "messages": [
            {
                "id": r[0],
                "role": r[1],
                "content": r[2],
                "created_at": r[3]
            } for r in rows
        ]
    }


def normalize_text(text: str) -> str:
  
    return unicodedata.normalize("NFKC", text)

def pdf_to_html_exact(pdf_bytes: bytes) -> str:
   
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    html = "<div class='pdf-container' style='position:relative;'>"

    for page_number, page in enumerate(doc, start=1):
       
        page_html = page.get_text("html")
       
        page_html = normalize_text(page_html)
       
        html += f"<div class='pdf-page' data-page='{page_number}'>{page_html}</div>"

    html += "</div>"
    return html
@app.post("/ocr")
async def ocr_endpoint(file: UploadFile = File(...)):
    pdf_bytes = await file.read()
    html = pdf_to_html_exact(pdf_bytes)
    return {"html": html}


@app.get("/conversations")
async def get_conversations(userid: str = Query(...)):
    try:
        cursor.execute(
            """
            SELECT id, user_id, title, created_at 
            FROM conversations
            WHERE user_id = %s
            ORDER BY created_at DESC
            """, 
            (userid,) 
        ) 
        
        rows = cursor.fetchall()

        conversations = [
            {
                "id": row[0],
                "userId": row[1],
                "title": row[2],
                "created_at": row[3]
            }
            for row in rows
        ]

        return {"conversations": conversations}

    except Exception as e:
        print(e)
        return {"error": "Failed to fetch conversations"}
    

@app.post("/ask")
async def ask(file: UploadFile = File(...), userid: str = Query(...)):
   
    conversation=upload_pdf(file,userid)
    file_extension = file.filename.split(".")[-1].lower()

    
    text = ""
    

    if file_extension in ["pdf"]:
       
        pdf = PdfReader(file.file)

        for page in pdf.pages:
            content = page.extract_text()
            if content:
                text += content + "\n"
    elif file_extension in ["ppt", "pptx"]:
        
        presentation = Presentation(file.file)

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    else:
        return {"error": "Unsupported file format. Use PDF or PPTX."}


    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = splitter.split_text(text)

    
    vectors = []
    ids = []

    for i, chunk in enumerate(chunks):
        try:
            vec = embeddings.embed_query(chunk)
            vectors.append(vec)
            ids.append(f"chunk-{i}")
            
         
                
        except Exception as e:
            print(e)
    

 
    index.upsert(
        vectors=[
            {
                "id": ids[i],
                "values": vectors[i],
                "metadata": {"text": chunks[i],"userId": userid  }
            }
            for i in range(len(chunks))
        ]
    )


    return conversation

def allowed_file(filename: str) -> bool:
    """Check if file has allowed extension"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def cleanup_file(filepath: Optional[str]) -> None:
    """Safely remove a file"""
    try:
        if filepath and os.path.exists(filepath):
            os.remove(filepath)
           
    except Exception as e:
        print(e)

@app.post("/convert-to-docx")
async def convert_to_docx(file: UploadFile = File(...)):
    """
    Convert uploaded PDF file to DOCX format
    
    Args:
        file: PDF file to convert
        
    Returns:
        FileResponse: Converted DOCX file
    """
    pdf_path = None
    docx_path = None
    
    try:
        # Validate file
        if not file:
          
            raise HTTPException(status_code=400, detail="No file provided")
        
        if not file.filename:
           
            raise HTTPException(status_code=400, detail="No filename provided")
        
       
        if not allowed_file(file.filename):
            
            raise HTTPException(
                status_code=400, 
                detail="Invalid file type. Only PDF files are allowed"
            )
        
   
        unique_id = str(uuid.uuid4())
        original_name = Path(file.filename).stem
        
        pdf_filename = f"{unique_id}_{original_name}.pdf"
        docx_filename = f"{unique_id}_{original_name}.docx"
        
        pdf_path = os.path.join(UPLOAD_FOLDER, pdf_filename)
        docx_path = os.path.join(UPLOAD_FOLDER, docx_filename)
        
        
        
        
        contents = await file.read()
        

        if len(contents) > MAX_FILE_SIZE:
            raise HTTPException(
                status_code=413, 
                detail=f"File too large. Maximum size is {MAX_FILE_SIZE / (1024*1024)}MB"
            )
        
        with open(pdf_path, 'wb') as f:
            f.write(contents)
        
        
        
   
        if not os.path.exists(pdf_path):
            raise HTTPException(status_code=500, detail="Failed to save PDF file")
        
       
        cv = Converter(pdf_path)
        cv.convert(docx_path, start=0, end=None)
        cv.close()
        
        
      
        if not os.path.exists(docx_path):
            raise HTTPException(status_code=500, detail="Failed to create DOCX file")
     
        return FileResponse(
            path=docx_path,
            media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            filename=f"{original_name}.docx",
            background=None  
        )
        
    except HTTPException:
     
        raise
        
    except Exception as e:
        
        raise HTTPException(
            status_code=500, 
            detail=f"Conversion failed: {str(e)}"
        )
    
    finally:
        # Clean up temporary files
        cleanup_file(pdf_path)
        # Note: We clean up DOCX after a delay to allow download
        # In production, consider using a background task or scheduled cleanup





if __name__ == "__main__":
    import uvicorn
    
  
    
    uvicorn.run(
        "app:app",  # Change "app" to your filename if different
        host="0.0.0.0",
        port=8000,
        reload=True,
        
    )